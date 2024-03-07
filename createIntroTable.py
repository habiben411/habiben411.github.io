#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Created on Tue Jan 23 06:47:13 2024

@author: hamidsarve
"""
arabic_and_phoenician_letters = [
    {'letter': 'ا', 'transliteration': 'alif', 'romanized': 'a', 'is_vowel': True, 'phoenician': '𐤀'},
    {'letter': 'ب', 'transliteration': 'ba', 'romanized': 'b', 'is_vowel': False, 'phoenician': '𐤁'},
    {'letter': 'ت', 'transliteration': 'ta', 'romanized': 't', 'is_vowel': False, 'phoenician': '𐤂'},
    {'letter': 'ث', 'transliteration': 'tha', 'romanized': 'th', 'is_vowel': False, 'phoenician': '𐤃'},
    {'letter': 'ج', 'transliteration': 'ja', 'romanized': 'j', 'is_vowel': False, 'phoenician': '𐤄'},
    {'letter': 'ح', 'transliteration': 'ha', 'romanized': 'h', 'is_vowel': False, 'phoenician': '𐤅'},
    {'letter': 'خ', 'transliteration': 'kha', 'romanized': 'kh', 'is_vowel': False, 'phoenician': '𐤆'},
    {'letter': 'د', 'transliteration': 'dal', 'romanized': 'd', 'is_vowel': False, 'phoenician': '𐤇'},
    {'letter': 'ذ', 'transliteration': 'dha', 'romanized': 'dh', 'is_vowel': False, 'phoenician': '𐤈'},
    {'letter': 'ر', 'transliteration': 'ra', 'romanized': 'r', 'is_vowel': False, 'phoenician': '𐤉'},
    {'letter': 'ز', 'transliteration': 'za', 'romanized': 'z', 'is_vowel': False, 'phoenician': '𐤊'},
    {'letter': 'س', 'transliteration': 'sa', 'romanized': 's', 'is_vowel': False, 'phoenician': '𐤋'},
    {'letter': 'ش', 'transliteration': 'sha', 'romanized': 'sh', 'is_vowel': False, 'phoenician': '𐤌'},
    {'letter': 'ص', 'transliteration': 'sad', 'romanized': 'ṣ', 'is_vowel': False, 'phoenician': '𐤍'},
    {'letter': 'ض', 'transliteration': 'dad', 'romanized': 'ḍ', 'is_vowel': False, 'phoenician': '𐤎'},
    {'letter': 'ط', 'transliteration': 'taa', 'romanized': 'ṭ', 'is_vowel': False, 'phoenician': '𐤏'},
    {'letter': 'ظ', 'transliteration': 'zaa', 'romanized': 'ẓ', 'is_vowel': False, 'phoenician': '𐤐'},
    {'letter': 'ع', 'transliteration': 'ain', 'romanized': 'ʿ', 'is_vowel': False, 'phoenician': '𐤑'},
    {'letter': 'غ', 'transliteration': 'ghain', 'romanized': 'gh', 'is_vowel': False, 'phoenician': '𐤒'},
    {'letter': 'ف', 'transliteration': 'fa', 'romanized': 'f', 'is_vowel': False, 'phoenician': '𐤓'},
    {'letter': 'ق', 'transliteration': 'qa', 'romanized': 'q', 'is_vowel': False, 'phoenician': '𐤔'},
    {'letter': 'ك', 'transliteration': 'kaf', 'romanized': 'k', 'is_vowel': False, 'phoenician': '𐤕'},
    {'letter': 'ل', 'transliteration': 'lam', 'romanized': 'l', 'is_vowel': False, 'phoenician': '𐤖'},
    {'letter': 'م', 'transliteration': 'mim', 'romanized': 'm', 'is_vowel': False, 'phoenician': '𐤗'},
    {'letter': 'ن', 'transliteration': 'nun', 'romanized': 'n', 'is_vowel': False, 'phoenician': '𐤘'},
    {'letter': 'ه', 'transliteration': 'haa', 'romanized': 'h', 'is_vowel': False, 'phoenician': '𐤙'},
    {'letter': 'و', 'transliteration': 'waw', 'romanized': 'w', 'is_vowel': True, 'phoenician': '𐤚'},
    {'letter': 'ي', 'transliteration': 'ya', 'romanized': 'y', 'is_vowel': True, 'phoenician': '𐤛'},
]
 
import random
 
swedish_words_from_arabic = [
    ('algebra', ('الجبر', 'al-jabr')),
    ('algoritm', ('الخوارزمي', 'al-Khwarizmi')),
    ('kaffe', ('القهوة', 'qahwa')),
    ('hasard', ('خطرة', 'hazārah')),
    ('soffa', ('صوفا', 'sufā')),
    ('siffra', ('صفر', 'ṣifr')),
    ('alkohol', ('الكحل', 'al-kuhl')),
    ('alkalisk', ('القلوي', 'al-qallī')),
    ('mascara', ('المسكرة', 'al-muskarrah')),
    ('azurit', ('اللازورد', 'al-lāzaward')),
    ('elixir', ('الإكسير', 'al-iksir')),
    ('safari', ('السفر', 'ṣafar')),
    ('socker', ('السكر', 'sukkar')),
    ('talisman', ('طِلَّسْم', 'tilismān')),
    ('zenit', ('سَمْت الرَأْس‎', 'samt ar-ras')),
    ('tariff', ('التعريفة', 'al-taʿrīfah')),
    ('admiral', ('أَمِير اَلبَحْر', 'amīr al-baḥr')),
    ('aprikos', ('البرقوق', 'al-baraqūq')),
    ('arsenal', ('دار الصناعة', 'dār al-ṣināʿa')),
    ('alkemi', ('الكيمياء', 'al-kimiyā')),
    ('madrass', ('المرتبة', 'al-murabbaṭah')),

]

random.shuffle(swedish_words_from_arabic)

swedish_words_from_arabic = [('average', (' عَوَار‎ ', 'awar'))]+  swedish_words_from_arabic 
output = ""
delay = 1250
tabs = '\t\t\t\t\t'
for (swe,ar) in swedish_words_from_arabic:
    text = ar[1] + "  " + ar[0]

    data_fade = 'data-autoslide=' + str(delay)
    output += tabs + 'section data-transition="fade" ' + data_fade + '>\n'
    output += tabs + '\t<div><h1>' + swe + '</h1>'
    
    output += '</div><br><br>'
    output += '<div><h3>' + text + '</h3>'
    output += '</div>\n'
    
    output +=  tabs + '</section>\n'


output = "".join(output.rsplit(data_fade, 1))


print(output)

if 0:
    from pptx import Presentation
    from pptx.util import Inches
    
    # Create a presentation object
    prs = Presentation()
    # Add a slide with a blank layout (usually layout 6)
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Define the table size and position
    x, y, cx, cy = Inches(2), Inches(2), Inches(4), Inches(3)
    # Add a table placeholder in the shape of a 4x7 table
    width = 7
    shape = slide.shapes.add_table(rows=4, cols=width, left=x, top=y, width=cx, height=cy)
    table = shape.table
    
    # Populate the table with Arabic and Phoenician letters
    for idx, letter_info in enumerate(arabic_and_phoenician_letters):
        
        row = idx // width  # Calculate the row index for each letter
        col = idx % width        # Calculate the column index for each letter
        
        print("idx " + str(idx) + " row " + str(row) + " " + "col " + str(col))
    
        # Set the Arabic letter in the first cell
        cell = table.cell(row, width-col-1)
        cell.text = letter_info['letter']
        cell.text = f"{letter_info['letter']}\n{letter_info['romanized']}\n{letter_info['phoenician']}"
    
    # Save the presentation
    prs.save('arabic_and_phoenician_letters.pptx')
